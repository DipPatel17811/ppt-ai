"""Validation.

Two distinct validators:

* ``validate_presentation`` - checks the *AST* before rendering (bullet
  counts, chart consistency, empty fields, structure);
* ``LayoutValidator`` - checks the *rendered* .pptx (no overlaps, margins
  respected, text fits its box, minimum font sizes, image/icon bounds).

Both return a :class:`Report` of human-readable issues.  Rendering is never
blocked by warnings; ``raise_on_error`` is available for strict pipelines.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Sequence

from pptx import Presentation as PptxPresentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

from config import (
    EMU_PER_PT,
    MAX_BULLETS_PER_LIST,
    MAX_CHART_CATEGORIES,
    MAX_CHART_SERIES,
    MAX_COMPARISON_POINTS,
    MAX_DASHBOARD_METRICS,
    MAX_HIERARCHY_NODES,
    MAX_PHASES,
    MAX_PIE_SLICES,
    MAX_STEPS,
    MAX_TIMELINE_ITEMS,
    MIN_FONT_SIZE,
)
from schema import (
    BulletsSlide,
    ChartKind,
    ComparisonSlide,
    ConclusionSlide,
    CycleSlide,
    DashboardSlide,
    HierarchySlide,
    ProcessSlide,
    RoadmapSlide,
    SwoSlide,
    TimelineSlide,
)
from utils import fonts as font_utils
from utils.geometry import Rect


@dataclass
class Report:
    """Result of a validation pass."""

    issues: List[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return not self.issues

    def error(self, msg: str) -> None:
        self.issues.append(msg)

    def extend(self, other: "Report") -> None:
        self.issues.extend(other.issues)

    def __str__(self) -> str:
        if not self.issues:
            return "OK"
        return "\n".join(f"- {i}" for i in self.issues)


# ---------------------------------------------------------------------------
# AST validation
# ---------------------------------------------------------------------------


def _validate_bullets(slide, report: Report) -> None:
    if len(slide.bullets) > MAX_BULLETS_PER_LIST:
        report.error(f"Slide {slide.type!r}: too many bullets ({len(slide.bullets)})")
    for i, item in enumerate(slide.bullets):
        if len(item.text) > 160:
            report.error(f"Slide {slide.type!r}: bullet {i} exceeds 160 chars")


def validate_presentation(ast) -> Report:
    """Validate an already-parsed Presentation AST and return a report."""
    report = Report()
    if not ast.title.strip():
        report.error("Presentation has an empty title")

    for idx, slide in enumerate(ast.slides):
        if isinstance(slide, BulletsSlide):
            _validate_bullets(slide, report)
        elif isinstance(slide, ComparisonSlide):
            if len(slide.left.points) > MAX_COMPARISON_POINTS:
                report.error(f"Slide {idx}: left column exceeds {MAX_COMPARISON_POINTS} points")
            if len(slide.right.points) > MAX_COMPARISON_POINTS:
                report.error(f"Slide {idx}: right column exceeds {MAX_COMPARISON_POINTS} points")
        elif isinstance(slide, TimelineSlide):
            if len(slide.items) > MAX_TIMELINE_ITEMS:
                report.error(f"Slide {idx}: too many timeline items ({len(slide.items)})")
            for i, item in enumerate(slide.items):
                if item.detail and len(item.detail) > 90:
                    report.error(f"Slide {idx}: timeline item {i} detail too long")
        elif isinstance(slide, ProcessSlide):
            if len(slide.steps) > MAX_STEPS:
                report.error(f"Slide {idx}: too many process steps ({len(slide.steps)})")
            for i, step in enumerate(slide.steps):
                if len(step) > 60:
                    report.error(f"Slide {idx}: process step {i} too long ({len(step)} chars)")
        elif isinstance(slide, RoadmapSlide):
            if len(slide.phases) > MAX_PHASES:
                report.error(f"Slide {idx}: too many roadmap phases ({len(slide.phases)})")
        elif isinstance(slide, CycleSlide):
            if len(slide.stages) > MAX_STEPS:
                report.error(f"Slide {idx}: too many cycle stages ({len(slide.stages)})")
        elif isinstance(slide, HierarchySlide):
            count = 0

            def walk(node):
                nonlocal count
                count += 1
                for ch in node.children:
                    walk(ch)

            walk(slide.root)
            if count > MAX_HIERARCHY_NODES:
                report.error(f"Slide {idx}: hierarchy too large ({count} nodes)")
        elif isinstance(slide, DashboardSlide):
            if len(slide.metrics) > MAX_DASHBOARD_METRICS:
                report.error(f"Slide {idx}: too many dashboard metrics ({len(slide.metrics)})")
        elif isinstance(slide, SwoSlide):
            for q in (slide.strengths, slide.weaknesses, slide.opportunities, slide.threats):
                if len(q.items) > 6:
                    report.error(f"Slide {idx}: SWOT quadrant {q.title!r} has too many items")
        elif isinstance(slide, ConclusionSlide):
            if len(slide.takeaways) > 6:
                report.error(f"Slide {idx}: too many takeaways")

        if hasattr(slide, "chart") and slide.chart is not None:
            _validate_chart(slide.chart, idx, report)
    return report


def _validate_chart(chart, idx: int, report: Report) -> None:
    if chart.kind == ChartKind.SCATTER:
        if not chart.scatter_points:
            report.error(f"Slide {idx}: scatter chart has no points")
        return
    if not chart.categories:
        report.error(f"Slide {idx}: chart has no categories")
        return
    if len(chart.categories) > MAX_CHART_CATEGORIES:
        report.error(f"Slide {idx}: too many chart categories ({len(chart.categories)})")
    if len(chart.series) > MAX_CHART_SERIES:
        report.error(f"Slide {idx}: too many chart series ({len(chart.series)})")
    for s in chart.series:
        if len(s.values) != len(chart.categories):
            report.error(f"Slide {idx}: series {s.name!r} length mismatch")
    if chart.kind == ChartKind.PIE and len(chart.categories) > MAX_PIE_SLICES:
        report.error(f"Slide {idx}: too many pie slices")


# ---------------------------------------------------------------------------
# Rendered .pptx validation
# ---------------------------------------------------------------------------


def _shape_rect(shape) -> Rect:
    try:
        x = shape.left / EMU_PER_PT
        y = shape.top / EMU_PER_PT
        w = shape.width / EMU_PER_PT
        h = shape.height / EMU_PER_PT
    except Exception:
        return Rect(0, 0, 0, 0)
    return Rect(x, y, w, h)


def _text_blocks(shape) -> List[tuple]:
    """Return [(font, size, bold, text)] for every run in the shape."""
    out = []
    try:
        tf = shape.text_frame
    except Exception:
        return out
    for para in tf.paragraphs:
        for run in para.runs:
            font = run.font
            try:
                size = font.size.pt if font.size else 14.0
            except Exception:
                size = 14.0
            try:
                bold = bool(font.bold)
            except Exception:
                bold = False
            family = font.name or "Segoe UI"
            out.append((family, size, bold, run.text))
    return out


class LayoutValidator:
    """Validates a rendered .pptx for overlap, margins, text fit, fonts."""

    def __init__(self, max_overlap_epsilon: float = 1.0,
                 text_overflow_tolerance: float = 1.15) -> None:
        self.epsilon = max_overlap_epsilon
        self.tolerance = text_overflow_tolerance

    def validate(self, path: str | Path) -> Report:
        report = Report()
        prs = PptxPresentation(str(path))
        sw = prs.slide_width / EMU_PER_PT
        sh = prs.slide_height / EMU_PER_PT
        page = Rect(0, 0, sw, sh)

        for si, slide in enumerate(prs.slides, start=1):
            self._validate_slide(slide, si, page, report)
        return report

    def _validate_slide(self, slide, si: int, page: Rect, report: Report) -> None:
        shapes = [s for s in slide.shapes]
        # out-of-bounds (skip full-bleed background/overlay rects + icon glyphs)
        for i, s in enumerate(shapes):
            r = _shape_rect(s)
            if r.width <= 0 or r.height <= 0:
                continue
            if self._is_background(r, page):
                continue
            try:
                if s.name.startswith("_icon_"):
                    continue
            except Exception:
                pass
            if not page.contains(r, epsilon=2.0):
                report.error(f"Slide {si}: shape {i} ({s.shape_type}) leaves the slide "
                             f"(x={r.x:.0f}, y={r.y:.0f})")
        # overlaps (top-level shapes only). Only "solid vs solid" and
        # "text vs text" count; a solid card is expected to sit under its own
        # text box, connectors/placeholders/lines are geometry not boxes, and
        # icon glyph pieces are tagged "_icon_*".
        for i in range(len(shapes)):
            for j in range(i + 1, len(shapes)):
                a, b = _shape_rect(shapes[i]), _shape_rect(shapes[j])
                if a.width <= 0 or b.width <= 0:
                    continue
                if (self._is_background(a, page) or self._is_background(b, page)):
                    continue
                ka, kb = self._box_kind(shapes[i]), self._box_kind(shapes[j])
                if ka == "none" or kb == "none":
                    continue
                if not (ka == kb):
                    continue
                if ka == "solid":
                    # nested badges/circle markers are expected composites,
                    # as are decorative ovals (process nodes + number badges)
                    if a.contains(b, epsilon=self.epsilon) or b.contains(a, epsilon=self.epsilon):
                        continue
                    if (shapes[i].auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL
                            and shapes[j].auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL):
                        continue
                inter = a.intersected(b)
                if inter is not None and min(a.area, b.area) > 0 \
                        and inter.area <= 0.15 * min(a.area, b.area):
                    continue
                if a.overlaps(b, epsilon=self.epsilon):
                    report.error(f"Slide {si}: shapes {i} and {j} overlap")
        # text fit + fonts
        for i, s in enumerate(shapes):
            self._check_text(s, i, si, report)

    @staticmethod
    def _is_background(r: Rect, page: Rect) -> bool:
        """A shape that covers most of the slide is a background/overlay."""
        return r.area > page.area * 0.90

    @staticmethod
    def _box_kind(shape) -> str:
        """Classify a shape for overlap purposes.

        * ``solid`` - has a visible fill (card/box/band);
        * ``text``  - a text box carrying real content;
        * ``none``  - connectors, placeholders, strokes-only freeforms,
          charts, pictures, groups and composed icon glyph pieces (tagged
          ``_icon_*``) - geometry rather than a box.
        """
        try:
            if shape.name.startswith("_icon_"):
                return "none"
        except Exception:
            pass
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            try:
                fill_type = shape.fill.type
            except Exception:
                fill_type = None
            if fill_type is None or fill_type == MSO_FILL_TYPE.BACKGROUND:
                return "text"
            return "solid"
        try:
            fill_type = shape.fill.type
        except Exception:
            fill_type = None
        if fill_type is None or fill_type == MSO_FILL_TYPE.BACKGROUND:
            return "none"
        return "solid"

    def _check_text(self, shape, i: int, si: int, report: Report) -> None:
        blocks = _text_blocks(shape)
        if not blocks:
            return
        rect = _shape_rect(shape)
        if rect.width <= 0:
            return
        try:
            tf = shape.text_frame
            wrap = bool(tf.word_wrap)
        except Exception:
            wrap = True
        total_h = 0.0
        seen_text = False
        for family, size, bold, text in blocks:
            if not text:
                continue
            seen_text = True
            if size < MIN_FONT_SIZE - 0.5:
                report.error(f"Slide {si}: shape {i} uses font size {size:.0f}pt (< {MIN_FONT_SIZE})")
            if wrap:
                lines = len(font_utils.wrap(text, family, size, rect.width, bold=bold))
            else:
                lines = max(1, len(text) // 80)
            total_h += lines * size * 1.2
        if seen_text and total_h > rect.height * self.tolerance:
            report.error(f"Slide {si}: shape {i} text may overflow "
                         f"(needs ~{total_h:.0f}pt, box {rect.height:.0f}pt)")
