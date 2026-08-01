"""The PowerPoint compiler.

This is the *only* module allowed to know about OOXML details.  It converts
the resolved, type-safe layout output into an editable python-pptx
Presentation: native shapes, native charts, real text runs, groups, images,
notes, transitions and animations.

Key invariants:

* every coordinate comes in as points and is converted to EMU here;
* charts are created as native PowerPoint charts (fully editable);
* images are placed with deterministic fit/fill cropping (never stretched
  out of proportion);
* animations/transitions/morph are injected as raw OOXML by the ``ppt``
  package -- this module is their gateway.
"""

from __future__ import annotations

import tempfile
import urllib.request
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Sequence, Tuple

from lxml import etree

from pptx import Presentation as PptxPresentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from config import ASPECT_RATIOS, EMU_PER_PT
from schema import ChartKind, ChartSpec, FitMode, ImageSpec, TextAlign, VerticalAlign
from utils.colors import normalize as _normalize_color
from utils.geometry import Rect


# ---------------------------------------------------------------------------
# Small typed structures passed from the renderer
# ---------------------------------------------------------------------------


@dataclass
class RunSpec:
    """One styled text run."""

    text: str
    font: Optional[str] = None
    size: Optional[float] = None
    color: Optional[str] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    letter_spacing: Optional[float] = None


@dataclass
class ParaSpec:
    """One paragraph (a sequence of runs plus paragraph-level styling)."""

    runs: List[RunSpec] = field(default_factory=list)
    align: TextAlign = TextAlign.LEFT
    space_before: float = 0.0
    space_after: float = 0.0
    line_spacing: Optional[float] = None
    bullet: bool = False
    bullet_char: str = "•"
    bullet_color: Optional[str] = None
    bullet_size_pct: float = 100.0
    indent_pt: Optional[float] = None


def _emu(value: float) -> Emu:
    return Emu(int(round(value * EMU_PER_PT)))


def _align(align: TextAlign) -> PP_ALIGN:
    return {
        TextAlign.LEFT: PP_ALIGN.LEFT,
        TextAlign.CENTER: PP_ALIGN.CENTER,
        TextAlign.RIGHT: PP_ALIGN.RIGHT,
    }[align]


def _anchor(anchor: VerticalAlign) -> MSO_ANCHOR:
    return {
        VerticalAlign.TOP: MSO_ANCHOR.TOP,
        VerticalAlign.MIDDLE: MSO_ANCHOR.MIDDLE,
        VerticalAlign.BOTTOM: MSO_ANCHOR.BOTTOM,
    }[anchor]


def _rgb(hex_color: str) -> RGBColor:
    h = _normalize_color(hex_color)
    return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))


def set_fill_alpha(shape, alpha: float) -> None:
    """Set solid-fill opacity to ``alpha`` in [0, 1] (0 = transparent)."""
    sp_pr = shape._element.spPr
    for el in sp_pr.iter():
        if el.tag == qn("a:srgbClr"):
            el.append(el.makeelement(qn("a:alpha"), {"val": str(int(round(alpha * 1000)))}))
            break


def set_arrow(shape, tail: bool = True, head: bool = False) -> None:
    """Add arrowheads to a line/connector shape."""
    ln = shape.line._get_or_add_ln()
    if tail:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if head:
        ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}))


class Compiler:
    """Deterministic bridge between the layout pipeline and python-pptx."""

    def __init__(self) -> None:
        self._prs: Optional[PptxPresentation] = None
        self._temp_files: List[str] = []

    # -- lifecycle -------------------------------------------------------
    def create(self, aspect: str = "16:9") -> None:
        w_pt, h_pt = ASPECT_RATIOS[aspect]
        self._prs = PptxPresentation()
        self._prs.slide_width = _emu(w_pt)
        self._prs.slide_height = _emu(h_pt)
        self._prs.slide_layouts  # ensure layouts loaded lazily

    @property
    def prs(self) -> PptxPresentation:
        if self._prs is None:
            raise RuntimeError("Compiler.create() must be called first")
        return self._prs

    def save(self, path: Path) -> None:
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        self._cleanup_temp()

    def _cleanup_temp(self) -> None:
        for p in self._temp_files:
            try:
                Path(p).unlink(missing_ok=True)
            except Exception:
                pass
        self._temp_files.clear()

    # -- slides ----------------------------------------------------------
    def add_slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def set_core_properties(self, title: str, author: str = "", subject: str = "") -> None:
        cp = self.prs.core_properties
        cp.title = title
        cp.author = author or "ppt-ai"
        cp.subject = subject

    def set_notes(self, slide, text: Optional[str]) -> None:
        if not text:
            return
        slide.notes_slide.notes_text_frame.text = text

    # -- shapes ----------------------------------------------------------
    def add_rect(self, slide, rect: Rect, fill: Optional[str] = None,
                 line: Optional[str] = None, line_width: float = 0.75,
                 radius: Optional[float] = None, alpha: float = 1.0) -> object:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, _emu(rect.x), _emu(rect.y),
                                       _emu(rect.width), _emu(rect.height))
        if radius is not None:
            try:
                shape.adjustments[0] = min(0.5, max(0.0, radius / min(rect.width, rect.height) * 0.5))
            except Exception:
                pass
        self._apply_fill(shape, fill, alpha)
        self._apply_line(shape, line, line_width)
        shape.shadow.inherit = False
        return shape

    def add_oval(self, slide, rect: Rect, fill: Optional[str] = None,
                 line: Optional[str] = None, line_width: float = 0.75,
                 alpha: float = 1.0) -> object:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, _emu(rect.x), _emu(rect.y),
                                       _emu(rect.width), _emu(rect.height))
        self._apply_fill(shape, fill, alpha)
        self._apply_line(shape, line, line_width)
        shape.shadow.inherit = False
        return shape

    def add_line(self, slide, x1: float, y1: float, x2: float, y2: float,
                 color: str, width: float = 1.5, dash: bool = False,
                 arrow: bool = False) -> object:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                          _emu(x1), _emu(y1), _emu(x2), _emu(y2))
        self._apply_line(conn, color, width)
        if dash:
            try:
                conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            except Exception:
                pass
        if arrow:
            set_arrow(conn, tail=True)
        return conn

    def add_freeform(self, slide, points: Sequence[Tuple[float, float]],
                     closed: bool = False, fill: Optional[str] = None,
                     line: Optional[str] = None, line_width: float = 1.0,
                     collection=None) -> object:
        """Draw a polygon/polyline from point-space coordinates.

        Coordinates are in points; ``scale=EMU_PER_PT`` converts them to
        OOXML EMU at the shape's bounding-box level (a single scaling, never
        a double conversion).
        """
        shapes = collection.shapes if collection is not None else slide.shapes
        if len(points) < 2:
            return None
        fb = shapes.build_freeform(points[0][0], points[0][1], scale=EMU_PER_PT)
        fb.add_line_segments([(x, y) for x, y in points[1:]], close=closed)
        shape = fb.convert_to_shape()
        self._apply_fill(shape, fill if closed else None, 1.0)
        self._apply_line(shape, line, line_width)
        shape.shadow.inherit = False
        return shape

    def add_textbox(self, slide, rect: Rect, paras: Sequence[ParaSpec],
                    anchor: VerticalAlign = VerticalAlign.TOP,
                    wrap: bool = True, margin: float = 4.0,
                    autofit: bool = True) -> object:
        tb = slide.shapes.add_textbox(_emu(rect.x), _emu(rect.y),
                                      _emu(rect.width), _emu(rect.height))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = _anchor(anchor)
        tf.auto_size = MSO_AUTO_SIZE.NONE if autofit else MSO_AUTO_SIZE.NONE
        tf.margin_left = tf.margin_right = _emu(margin)
        tf.margin_top = tf.margin_bottom = _emu(margin)
        self._fill_text_frame(tf, paras)
        return tb

    def _fill_text_frame(self, tf, paras: Sequence[ParaSpec]) -> None:
        first = True
        for para in paras:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = _align(para.align)
            p.space_before = Pt(para.space_before)
            p.space_after = Pt(para.space_after)
            if para.line_spacing is not None:
                p.line_spacing = para.line_spacing
            if para.indent_pt is not None:
                self._set_indent(p, para.indent_pt)
            for run in para.runs:
                r = p.add_run()
                r.text = run.text
                f = r.font
                f.name = run.font
                if run.size is not None:
                    f.size = Pt(run.size)
                if run.color is not None:
                    f.color.rgb = _rgb(run.color)
                if run.bold is not None:
                    f.bold = run.bold
                if run.italic is not None:
                    f.italic = run.italic
                if run.letter_spacing is not None:
                    self._set_letter_spacing(r, run.letter_spacing)
                self._set_font_hint(r, run.font)
            if para.bullet:
                self._set_bullet(p, para.bullet_char, para.bullet_color, para.bullet_size_pct)

    def _set_indent(self, paragraph, indent_pt: float) -> None:
        pPr = paragraph._p.get_or_add_pPr()
        pPr.set("marL", str(int(indent_pt * 100)))
        pPr.set("indent", str(-int(indent_pt * 100)))

    def _set_bullet(self, paragraph, char: str, color: Optional[str], size_pct: float) -> None:
        pPr = paragraph._p.get_or_add_pPr()
        # remove any existing bullet nodes
        for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buClr", "a:buFont", "a:buSzPct", "a:buSzPts"):
            for el in pPr.findall(qn(tag)):
                pPr.remove(el)
        if color:
            buClr = etree.SubElement(pPr, qn("a:buClr"))
            srgb = etree.SubElement(buClr, qn("a:srgbClr"))
            srgb.set("val", _normalize_color(color)[1:])
        buSz = etree.SubElement(pPr, qn("a:buSzPct"))
        buSz.set("val", str(int(size_pct * 1000)))
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "Arial")
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", char)

    def _set_letter_spacing(self, run, spc_pt: float) -> None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(spc_pt * 100)))

    def _set_font_hint(self, run, font: Optional[str]) -> None:
        if not font:
            return
        rPr = run._r.get_or_add_rPr()
        ea = rPr.makeelement(qn("a:ea"), {"typeface": font})
        latin = rPr.find(qn("a:latin"))
        if latin is not None:
            rPr.insert(list(rPr).index(latin) + 1, ea)
        else:
            rPr.append(ea)

    def _apply_fill(self, shape, color: Optional[str], alpha: float) -> None:
        if color is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(color)
            if alpha < 0.999:
                set_fill_alpha(shape, alpha)

    def _apply_line(self, shape, color: Optional[str], width: float) -> None:
        if color is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = _rgb(color)
            shape.line.width = Pt(width)

    # -- groups ----------------------------------------------------------
    def add_group(self, slide, local_rect: Rect, left: float, top: float):
        """Create a group positioned at (left, top) whose children use local
        coordinates within ``local_rect``."""
        group = slide.shapes.add_group_shape()
        group.left = _emu(left)
        group.top = _emu(top)
        group.width = _emu(local_rect.width)
        group.height = _emu(local_rect.height)
        return group

    # -- images ----------------------------------------------------------
    def add_picture(self, slide, rect: Rect, spec: ImageSpec) -> object:
        path = self._resolve_image(spec.source)
        try:
            from PIL import Image
        except Exception:
            Image = None
        if Image is not None:
            with Image.open(path) as im:
                iw, ih = im.size
        else:
            iw = ih = 100
        tw, th = rect.width, rect.height

        if spec.fit == FitMode.CONTAIN:
            scale = min(tw / iw, th / ih)
            pw, ph = iw * scale, ih * scale
            left = rect.x + (tw - pw) * spec.align_x
            top = rect.y + (th - ph) * spec.align_y
            picture = slide.shapes.add_picture(str(path), _emu(left), _emu(top),
                                               _emu(pw), _emu(ph))
        else:  # COVER
            scale = max(tw / iw, th / ih)
            pw, ph = iw * scale, ih * scale
            left = rect.x - (pw - tw) * spec.align_x
            top = rect.y - (ph - th) * spec.align_y
            picture = slide.shapes.add_picture(str(path), _emu(left), _emu(top),
                                               _emu(pw), _emu(ph))
            if pw > tw:
                fh = (pw - tw) / pw
                picture.crop_left = spec.align_x * fh
                picture.crop_right = (1.0 - spec.align_x) * fh
            if ph > th:
                fv = (ph - th) / ph
                picture.crop_top = spec.align_y * fv
                picture.crop_bottom = (1.0 - spec.align_y) * fv
        if spec.alt_text:
            picture.shape_id  # ensure element exists
            picture._element.set("descr", spec.alt_text)
        return picture

    def _resolve_image(self, source: str) -> Path:
        if source.lower().startswith(("http://", "https://")):
            suffix = Path(source).suffix or ".png"
            tmp = tempfile.mkstemp(suffix=suffix, prefix="ppt_ai_img_")[1]
            urllib.request.urlretrieve(source, tmp)  # noqa: S310
            self._temp_files.append(tmp)
            return Path(tmp)
        path = Path(source)
        if not path.exists():
            raise FileNotFoundError(f"Image source not found: {source}")
        return path

    # -- charts ----------------------------------------------------------
    def add_chart(self, slide, rect: Rect, spec: ChartSpec) -> object:
        chart_type = self._chart_type(spec)
        if spec.kind == ChartKind.SCATTER:
            chart_data = XyChartData()
            series = chart_data.add_series(spec.series[0].name if spec.series else "Series 1")
            for x, y in (spec.scatter_points or []):
                series.add_data_point(x, y)
        else:
            chart_data = CategoryChartData()
            chart_data.categories = list(spec.categories)
            for s in spec.series:
                chart_data.add_series(s.name, tuple(s.values))

        gframe = slide.shapes.add_chart(chart_type, _emu(rect.x), _emu(rect.y),
                                        _emu(rect.width), _emu(rect.height), chart_data)
        chart = gframe.chart
        self._style_chart(chart, spec)
        return chart

    def _chart_type(self, spec: ChartSpec):
        if spec.kind == ChartKind.PIE:
            return XL_CHART_TYPE.PIE
        if spec.kind == ChartKind.LINE:
            return XL_CHART_TYPE.LINE_MARKERS
        if spec.kind == ChartKind.AREA:
            return XL_CHART_TYPE.AREA_STACKED if spec.stacked else XL_CHART_TYPE.AREA
        if spec.kind == ChartKind.SCATTER:
            return XL_CHART_TYPE.SCATTER_MARKERS
        if spec.horizontal:
            return XL_CHART_TYPE.BAR_STACKED if spec.stacked else XL_CHART_TYPE.BAR_CLUSTERED
        return XL_CHART_TYPE.COLUMN_STACKED if spec.stacked else XL_CHART_TYPE.COLUMN_CLUSTERED

    def _style_chart(self, chart, spec: ChartSpec) -> None:
        try:
            chart.has_title = bool(spec.title)
            if spec.title:
                chart.chart_title.text_frame.text = spec.title
                chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
                chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
        except Exception:
            pass
        try:
            chart.has_legend = spec.show_legend
            if spec.show_legend:
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                chart.legend.font.size = Pt(10)
        except Exception:
            pass
        try:
            plot = chart.plots[0]
            if spec.show_value_labels:
                plot.has_data_labels = True
                labels = plot.data_labels
                labels.show_value = True
                labels.show_category_name = False
                labels.show_series_name = False
                labels.number_format = "0.0"
                labels.number_format_is_linked = False
                labels.font.size = Pt(10)
        except Exception:
            pass
        # Series colours.
        try:
            palette = [s.color for s in spec.series] or [None] * len(chart.series)
            for i, series in enumerate(chart.series):
                color = palette[i] if i < len(palette) else None
                if color is None:
                    continue
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = _rgb(color)
                if spec.kind in (ChartKind.LINE, ChartKind.SCATTER):
                    series.format.line.color.rgb = _rgb(color)
                    series.format.line.width = Pt(2)
                    series.marker.format.fill.fore_color.rgb = _rgb(color)
        except Exception:
            pass
        # Axes.
        if spec.kind not in (ChartKind.PIE,):
            try:
                va = chart.value_axis
                va.has_major_gridlines = spec.show_gridlines
                va.tick_labels.font.size = Pt(10)
                if spec.y_label:
                    va.has_title = True
                    va.axis_title.text_frame.text = spec.y_label
            except Exception:
                pass
            try:
                ca = chart.category_axis
                ca.tick_labels.font.size = Pt(10)
                if spec.x_label:
                    ca.has_title = True
                    ca.axis_title.text_frame.text = spec.x_label
            except Exception:
                pass

    # -- slide XML hooks (delegated to the ppt package) ------------------
    @staticmethod
    def slide_element(slide) -> etree._Element:
        return slide._element

    @staticmethod
    def shape_ids(slide) -> List[int]:
        return [sh.shape_id for sh in slide.shapes]
